#!/usr/bin/env python3
"""ACMS report service — turns live ACMS data into PDF and PPTX deliverables.

Every number in a generated report is fetched from the ACMS API and totalled
here in Python. The agent that calls this service may pass a `narrative` — free
prose for the commentary block — but it never supplies a figure. That split is
deliberate: asked to add up a BOQ, the model produced six correct line totals
and a grand total a thousand dollars out.

Endpoints
    GET  /health
    POST /generate   {report, format, opportunityId?, days?, narrative?}
    GET  /files/<name>

PDF is rendered by the Chrome that is already installed on this box — nothing
else on the machine shapes Arabic correctly. PPTX is built with python-pptx.

Run:  python3 service.py        (PORT env var, default 3025)
"""

import base64
import hashlib
import html as html_mod
import json
import os
import random
import re
import shutil
import subprocess
import tempfile
import threading
import time
import urllib.error
import urllib.parse
import urllib.request
import uuid
from datetime import datetime, timedelta, timezone
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer

API = os.environ.get("ACMS_API", "http://100.122.6.64:4010/api")
EMAIL = os.environ.get("ACMS_EMAIL", "ai.agent@afro.example")
PASSWORD = os.environ.get("ACMS_PASSWORD", "AcmsAgent#n8n2026")
PORT = int(os.environ.get("PORT", "3025"))
# One interface, not all of them. Default is the tailscale address: the n8n
# container and the users' browsers both reach the service there, and nothing
# legitimate arrives over the office LAN.
BIND = os.environ.get("ACMS_BIND", "100.122.6.64")
PUBLIC_BASE = os.environ.get("REPORTS_BASE", "http://100.122.6.64:%d" % PORT)

HERE = os.path.dirname(os.path.abspath(__file__))
FILES = os.path.join(HERE, "files")
os.makedirs(FILES, exist_ok=True)

CHROME = shutil.which("google-chrome") or shutil.which("google-chrome-stable")

SAFE_NAME = re.compile(r"^[A-Za-z0-9._-]+$")

# ------------------------------------------------- pending approvals store
#
# The chat approval gate used to keep proposals in n8n's workflow static data.
# That turned out not to survive between executions here — the proposal was
# written, the confirmation could not find it, and a confirmed action silently
# did nothing. Keeping the state in this service instead makes it explicit,
# inspectable, and independent of how n8n runs Code nodes.

PENDING_TTL = 10 * 60
_pending = {}
_pending_lock = threading.Lock()


def pending_put(session_id, proposal):
    with _pending_lock:
        _prune_pending()
        proposal["ts"] = time.time()
        _pending[session_id] = proposal


def pending_get(session_id):
    with _pending_lock:
        _prune_pending()
        return _pending.get(session_id)


def pending_pop(session_id):
    with _pending_lock:
        _prune_pending()
        return _pending.pop(session_id, None)


def _prune_pending():
    now = time.time()
    for k in [k for k, v in _pending.items() if now - v.get("ts", 0) > PENDING_TTL]:
        del _pending[k]


# ----------------------------------------------------------- action whitelist
#
# The action agent used to emit an @@ACTION{...}@@ marker in its answer and a
# Code node parsed it. That relied on the orchestrator relaying the specialist's
# text verbatim, and it did not: it paraphrased the sentence and dropped the
# marker, so a confirmed-looking request stored nothing. Proposing is now a real
# tool call into this endpoint — a side effect that no amount of rewording can
# lose — and the whitelist lives here, where the model cannot reach it.

UUID_RE = re.compile(
    r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$", re.I)

ID_FIELDS = ("opportunityId", "accountId", "approvalId", "discountId", "contactId")

OPP_UPDATABLE = ("scopeSummary", "solutionStrategy", "nextStep", "estimatedValue",
                 "proposedPrice", "probability", "forecastCategory", "health",
                 "expectedCloseDate")


def _opt(dst, src, keys):
    for k in keys:
        v = src.get(k)
        if v not in (None, ""):
            dst[k] = v
    return dst


ACTIONS = {
    "stage": dict(
        label="تغيير مرحلة فرصة", need=["opportunityId", "toStage"],
        url=lambda p: "/opportunities/%s/stage" % p["opportunityId"],
        body=lambda p: {"toStage": p["toStage"],
                        "reason": p.get("reason") or "عبر المساعد الذكي"}),
    "update": dict(
        label="تحديث بيانات فرصة", need=["opportunityId"], method="PATCH",
        url=lambda p: "/opportunities/%s" % p["opportunityId"],
        body=lambda p: _opt({}, p.get("fields") or p, OPP_UPDATABLE)),
    "status": dict(
        label="تغيير حالة فرصة", need=["opportunityId", "status"],
        url=lambda p: "/opportunities/%s/status" % p["opportunityId"],
        body=lambda p: _opt({"status": p["status"]}, p, ("exitReason", "exitNotes"))),
    "activity": dict(
        label="تسجيل نشاط", need=["type", "subject"],
        url=lambda p: "/activities",
        body=lambda p: _opt({"type": p["type"], "subject": p["subject"]}, p,
                            ("opportunityId", "accountId", "body", "dueAt",
                             "completed"))),
    "account_create": dict(
        label="إنشاء حساب/عميل جديد", need=["legalName", "type", "country"],
        url=lambda p: "/accounts",
        body=lambda p: _opt({"legalName": p["legalName"], "type": p["type"],
                             "country": p["country"]}, p,
                            ("tradeName", "industry", "city", "address", "website",
                             "taxId", "creditStatus", "paymentTermDays"))),
    "contact_create": dict(
        label="إضافة جهة اتصال", need=["accountId", "fullName"],
        url=lambda p: "/contacts",
        body=lambda p: _contact_body(p)),
    "opportunity_create": dict(
        label="إنشاء فرصة جديدة", need=["name", "accountId", "country"],
        url=lambda p: "/opportunities",
        body=lambda p: _opportunity_body(p)),
    "approval_decide": dict(
        label="البتّ في طلب موافقة", need=["approvalId", "decision"],
        url=lambda p: "/approvals/%s/decide" % p["approvalId"],
        body=lambda p: {"decision": p["decision"],
                        "note": p.get("note") or "عبر المساعد الذكي"}),
    "discount_decide": dict(
        label="البتّ في طلب خصم", need=["discountId", "approve"],
        url=lambda p: "/discounts/%s/decide" % p["discountId"],
        body=lambda p: {"approve": bool(p["approve"]),
                        "note": p.get("note") or "عبر المساعد الذكي"}),
}


def _contact_body(p):
    b = _opt({"accountId": p["accountId"], "fullName": p["fullName"]}, p,
             ("jobTitle", "email", "phone", "mobile", "influence", "notes"))
    if isinstance(p.get("isPrimary"), bool):
        b["isPrimary"] = p["isPrimary"]
    roles = p.get("roles")
    if roles:
        b["roles"] = roles if isinstance(roles, list) else [roles]
    return b


def _opportunity_body(p):
    b = _opt({"name": p["name"], "accountId": p["accountId"],
              "country": p["country"]}, p,
             ("industry", "source", "currency", "nextStep"))
    # A model that does not know the value sends 0, which reads as "worth
    # nothing" rather than "not priced yet". Leave it unset instead.
    try:
        v = float(p.get("estimatedValue"))
        if v > 0:
            b["estimatedValue"] = v
    except (TypeError, ValueError):
        pass
    return b


# Creating something that already exists is nearly always a retry of a
# half-finished attempt — that is how a zero-value opportunity ended up beside
# the real one.
DUP_CHECK = {
    "opportunity_create": ("/opportunities", "name", "فرصة"),
    "account_create": ("/accounts", "legalName", "حساب"),
    "contact_create": ("/contacts", "fullName", "جهة اتصال"),
}


def _norm(s):
    return " ".join(str(s or "").lower().split())


ERR_NO_RECORD = ("لا يوجد %s بالمعرّف %s في النظام. ابحث عن السجل الصحيح "
                 "أولًا ولا تستخدم معرّف سجل آخر.")

ERR_NAME_MISSING = ("ناقص %s: اذكر اسم %s كما طلبه المستخدم حتى يتحقّق النظام "
                    "من أن المعرّف يخصّ السجل نفسه.")

ERR_NAME_MISMATCH = ("المعرّف المُرسَل يخصّ %s «%s» (%s)، وليس «%s». لم يُقترح أي "
                     "إجراء. إن كان المطلوب غير موجود فقل ذلك صراحةً ولا "
                     "تستبدل به سجلًا آخر.")

# ---------------------------------------------------- reference verification
#
# A UUID that parses is not a UUID that means anything. The case that prompted
# this: asked to create an opportunity for "Zzz Telecom" — an account that does
# not exist — the agent supplied Vodafone's real accountId. Well formed,
# existent, and the wrong company. Format validation cannot catch that. Only
# resolving the id against ACMS and comparing it with the entity the model says
# it is acting on can.
#
# Only collections with a proven GET /{collection}/{id} are listed here.
# /approvals and /discounts have no such route on this build, so their ids stay
# format-checked only — fail open and documented rather than guessed at.
REF_LOOKUP = {
    "opportunityId": ("/opportunities", ("name",), "الفرصة"),
    "accountId": ("/accounts", ("legalName", "tradeName"), "الحساب"),
    "contactId": ("/contacts", ("fullName",), "جهة الاتصال"),
}

# Creates where the user names a company out loud are where a wrong id does the
# most damage and is least visible in the confirmation, so for these the model
# must also state the name it believes it is acting on, and the two must agree.
NAME_CONFIRM = {
    "opportunity_create": ("accountId", "accountName"),
    "contact_create": ("accountId", "accountName"),
}


def _ref_name(rec, fields):
    for f in fields:
        if rec.get(f):
            return str(rec[f])
    return ""


def _resolve_refs(params):
    """Resolve every referenced id to a real record. Returns (refs, error)."""
    refs = {}
    for field, (path, name_fields, label) in REF_LOOKUP.items():
        rid = params.get(field)
        if not rid:
            continue
        try:
            rec = acms.get("%s/%s" % (path, rid))
        except RuntimeError as e:
            if "404" in str(e):
                return None, ERR_NO_RECORD % (label, rid)
            # A lookup that failed for any other reason (network, 5xx) must not
            # block a legitimate proposal — the rule DUP_CHECK already follows.
            continue
        except Exception:
            continue
        if not rec:
            continue
        refs[field] = {"id": rid, "code": rec.get("code") or "",
                       "name": _ref_name(rec, name_fields), "label": label}
    return refs, None


def _same_entity(claimed, ref):
    """Does the name the model stated match the record its id resolved to?"""
    c = _norm(claimed)
    if not c:
        return False
    for candidate in (ref.get("name"), ref.get("code")):
        v = _norm(candidate)
        if v and (v == c or v in c or c in v):
            return True
    return False


def propose(session_id, action, params, summary=""):
    spec = ACTIONS.get(action)
    if not spec:
        return {"ok": False,
                "error": 'الإجراء "%s" غير مسموح به. المسموح: %s'
                         % (action, "، ".join(sorted(ACTIONS)))}

    missing = [k for k in spec["need"]
               if params.get(k) in (None, "", [])]
    if missing:
        return {"ok": False, "error": "بيانات ناقصة: " + "، ".join(missing)}

    for k in ID_FIELDS:
        v = params.get(k)
        if v and not UUID_RE.match(str(v)):
            return {"ok": False,
                    "error": 'القيمة "%s" ليست معرّف %s صالحًا. استخدم UUID من '
                             'أداة البحث لا الكود الظاهر للمستخدم.' % (v, k)}

    # Existence and identity, not just shape. This is the check the UUID regex
    # above cannot make: it rejects an id that names nothing, and an id that
    # names the wrong record.
    refs, ref_err = _resolve_refs(params)
    if ref_err:
        return {"ok": False, "error": ref_err}

    pair = NAME_CONFIRM.get(action)
    if pair:
        id_field, name_field = pair
        ref = refs.get(id_field)
        claimed = str(params.get(name_field) or "").strip()
        label = REF_LOOKUP[id_field][2]
        if not claimed:
            return {"ok": False, "error": ERR_NAME_MISSING % (name_field, label)}
        if ref and not _same_entity(claimed, ref):
            return {"ok": False,
                    "error": ERR_NAME_MISMATCH % (label, ref["name"],
                                                  ref["code"] or ref["id"],
                                                  claimed)}

    dup = DUP_CHECK.get(action)
    if dup:
        path, field, label = dup
        term = str(params.get(field) or "").strip()
        if len(term) >= 3:
            try:
                found = acms.get("%s?search=%s" % (path, urllib.parse.quote(term)))
                items = (found or {}).get("items", [])
            except Exception:
                items = []      # a failed lookup must not block a real create
            hit = next((x for x in items if _norm(x.get(field)) == _norm(term)), None)
            if hit:
                return {"ok": False,
                        "error": "يوجد %s بهذا الاسم بالفعل: %s — %s. لم يُقترح أي "
                                 "إنشاء؛ وضّح اسمًا مختلفًا أو اطلب تعديل الموجود."
                                 % (label, hit.get("code") or hit.get("id"),
                                    hit.get(field) or "")}

    body = spec["body"](params)
    if not body:
        return {"ok": False, "error": "لا توجد حقول صالحة للتحديث."}

    code = "%04d" % random.randint(1000, 9999)
    pending_put(session_id, {
        "code": code,
        "label": spec["label"],
        "summary": summary or spec["label"],
        # Method and path come from this table, never from the model.
        "method": spec.get("method", "POST"),
        "url": API + spec["url"](params),
        "body": body,
        "announced": False,
        # What the ids actually resolve to, so the confirmation can name the
        # target instead of repeating the model's claim about it.
        "resolved": refs,
    })
    return {"ok": True, "code": code, "label": spec["label"], "fields": body,
            "resolved": refs}


# --------------------------------------------------------------------- ACMS

class Acms:
    """Minimal ACMS client. The access token lives 15 minutes, so it is
    refetched whenever it is close to expiry rather than cached forever."""

    def __init__(self):
        self._token = None
        self._got_at = 0.0
        self._lock = threading.Lock()

    def _login(self):
        body = json.dumps({"email": EMAIL, "password": PASSWORD}).encode()
        req = urllib.request.Request(API + "/auth/login", data=body, method="POST")
        req.add_header("Content-Type", "application/json")
        with urllib.request.urlopen(req, timeout=30) as r:
            self._token = json.load(r)["accessToken"]
        self._got_at = time.time()

    def token(self):
        with self._lock:
            if self._token is None or time.time() - self._got_at > 10 * 60:
                self._login()
            return self._token

    def get(self, path):
        req = urllib.request.Request(API + path)
        req.add_header("Authorization", "Bearer " + self.token())
        try:
            with urllib.request.urlopen(req, timeout=60) as r:
                raw = r.read().decode()
            return json.loads(raw) if raw else None
        except urllib.error.HTTPError as e:
            raise RuntimeError("%s -> %s %s" % (path, e.code, e.read().decode()[:200]))


acms = Acms()


def num(v, default=0.0):
    try:
        return float(v)
    except (TypeError, ValueError):
        return default


def money(v, currency="USD"):
    return "{:,.0f} {}".format(num(v), currency)


def pct(v):
    return "—" if v is None else "{:,.2f}%".format(num(v))


# ------------------------------------------------------------ data gathering

METRIC_LABELS = {
    "WEIGHTED_PIPELINE": "الـPipeline المرجَّح",
    "WIN_RATE": "نسبة الفوز",
    "FORECAST_ACCURACY": "دقة التوقع",
    "OPEN_APPROVALS": "موافقات مفتوحة",
    "APPROVAL_WAIT": "متوسط انتظار الموافقة",
    "GROSS_MARGIN": "الهامش الإجمالي",
}

STAGE_LABELS = {
    "LEAD_INTAKE": "استقبال",
    "LEAD_QUALIFICATION": "تأهيل مبدئي",
    "OPPORTUNITY_QUALIFICATION": "تأهيل الفرصة",
    "SCOPE_DISCOVERY": "تحديد النطاق",
    "BID_STRATEGY_SOLUTION": "استراتيجية العطاء",
    "COSTING_SOURCING": "التكلفة والتوريد",
    "OPERATIONAL_FINANCIAL_REVIEW": "مراجعة تشغيلية ومالية",
    "MANAGEMENT_APPROVAL": "اعتماد الإدارة",
    "PROPOSAL_SUBMISSION": "تقديم العرض",
    "CLARIFICATIONS_NEGOTIATION": "استيضاحات وتفاوض",
    "AWARD_CONTRACTING": "الترسية والتعاقد",
    "PROJECT_HANDOVER": "تسليم للمشروع",
    "ACTUAL_PERFORMANCE_FEEDBACK": "الأداء الفعلي",
}


def gather_summary(days=60):
    """Everything the executive reports need, in one pass."""
    metrics = acms.get("/metrics/dashboard")
    active = acms.get("/opportunities?status=ACTIVE")["items"]
    lost = acms.get("/opportunities?status=LOST")["items"]
    closed = acms.get("/opportunities?status=CLOSED")["items"]
    deadlines = acms.get("/bids/deadlines?days=%d" % int(days)) or []
    approvals = acms.get("/approvals/my-queue") or []

    # Totals are computed here, never asked of a model.
    open_value = sum(num(o.get("estimatedValue")) for o in active)
    weighted = sum(num(o.get("estimatedValue")) * num(o.get("probability")) / 100.0
                   for o in active)
    won_value = sum(num(o.get("estimatedValue")) for o in closed)
    lost_value = sum(num(o.get("estimatedValue")) for o in lost)

    by_stage = {}
    for o in active:
        s = o.get("stage") or "—"
        e = by_stage.setdefault(s, {"count": 0, "value": 0.0})
        e["count"] += 1
        e["value"] += num(o.get("estimatedValue"))

    return {
        "asOf": datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC"),
        "metrics": metrics.get("metrics", []) if metrics else [],
        "active": sorted(active, key=lambda o: -num(o.get("estimatedValue"))),
        "lost": lost,
        "closed": closed,
        "deadlines": deadlines,
        "approvals": approvals,
        "openCount": len(active),
        "openValue": open_value,
        "weighted": weighted,
        "wonCount": len(closed),
        "wonValue": won_value,
        "lostCount": len(lost),
        "lostValue": lost_value,
        "byStage": sorted(by_stage.items(), key=lambda kv: -kv[1]["value"]),
        "days": int(days),
    }


def gather_cost_sheet(opportunity_id):
    opp = acms.get("/opportunities/%s" % opportunity_id)
    scenarios = acms.get("/opportunities/%s/costing" % opportunity_id) or []
    chosen = next((s for s in scenarios if s.get("isSelected")), None) or (
        scenarios[0] if scenarios else None)
    packages = []
    version = None
    if chosen and chosen.get("versions"):
        version = chosen["versions"][0]
        detail = acms.get("/costing/versions/%s" % version["id"])
        packages = (detail or {}).get("packages", []) or []

    rows = []
    grand_cost = grand_price = 0.0
    for pkg in packages:
        for item in pkg.get("items", []) or []:
            qty = num(item.get("quantity"))
            sell = num(item.get("sellingRate"))
            cost = sum(num(b.get("quantity")) * num(b.get("unitCost"))
                       for b in item.get("breakdown", []) or [])
            line_price = qty * sell
            grand_price += line_price
            grand_cost += cost
            rows.append({
                "package": pkg.get("name", "—"),
                "description": item.get("description", "—"),
                "unit": item.get("unit") or "",
                "quantity": qty,
                "sellingRate": sell,
                "linePrice": line_price,
                "lineCost": cost,
            })

    margin = ((grand_price - grand_cost) / grand_price * 100.0) if grand_price else None
    return {
        "asOf": datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC"),
        "opportunity": opp or {},
        "scenario": chosen or {},
        "version": version or {},
        "rows": rows,
        "grandCost": grand_cost,
        "grandPrice": grand_price,
        # Margin is on the selling price, never on cost — the system's own rule.
        "margin": margin,
        "currency": (chosen or {}).get("currency") or (opp or {}).get("currency") or "USD",
    }


# --------------------------------------------------------------- HTML → PDF

CSS = """
@page { size: A4 portrait; margin: 14mm; }
* { box-sizing: border-box; -webkit-print-color-adjust: exact; print-color-adjust: exact; }
body {
  margin: 0; direction: rtl; color: #10181c; background: #fff;
  font-family: "Segoe UI","Noto Kufi Arabic","Cairo",Tahoma,sans-serif;
  font-size: 12px; line-height: 1.65;
}
.mono { font-family: "Cascadia Mono",Consolas,"DejaVu Sans Mono",monospace; direction: ltr;
        font-variant-numeric: tabular-nums; }
header { border-bottom: 2px solid #0d6f7c; padding-bottom: 10px; margin-bottom: 18px; }
.eyebrow { font-family: "Cascadia Mono",Consolas,monospace; font-size: 9.5px;
           letter-spacing: .16em; color: #0d6f7c; direction: ltr; text-align: right; }
h1 { margin: 4px 0 2px; font-size: 22px; letter-spacing: -.01em; }
.sub { color: #52626b; font-size: 12px; }
h2 { font-size: 11px; font-family: "Cascadia Mono",Consolas,monospace; letter-spacing: .12em;
     color: #6f8089; text-transform: uppercase; direction: ltr; text-align: right;
     border-bottom: 1px solid #ccd6da; padding-bottom: 5px; margin: 22px 0 10px; }
.kpis { display: grid; grid-template-columns: repeat(3,1fr); gap: 8px; }
.kpi { border: 1px solid #ccd6da; border-radius: 3px; padding: 9px 11px; }
.kpi .k { font-size: 10.5px; color: #6f8089; }
.kpi .v { font-size: 17px; font-weight: 700; font-family: "Cascadia Mono",Consolas,monospace;
          direction: ltr; text-align: right; font-variant-numeric: tabular-nums; }
.kpi .n { font-size: 9.5px; color: #8b9aa2; }
table { width: 100%; border-collapse: collapse; font-size: 11px; }
th, td { text-align: right; padding: 6px 8px; border-bottom: 1px solid #dde5e8; vertical-align: top; }
th { font-size: 9.5px; color: #6f8089; text-transform: uppercase; letter-spacing: .06em;
     font-family: "Cascadia Mono",Consolas,monospace; direction: ltr; white-space: nowrap; }
td.n { font-family: "Cascadia Mono",Consolas,monospace; direction: ltr; text-align: right;
       font-variant-numeric: tabular-nums; white-space: nowrap; }
tr.total td { border-top: 2px solid #0d6f7c; border-bottom: none; font-weight: 700;
              background: #eef5f6; }
.narrative { border-inline-start: 3px solid #0d6f7c; padding: 2px 12px; color: #33454d;
             background: #f5f9fa; }
.empty { color: #8b9aa2; font-style: italic; }
footer { margin-top: 26px; padding-top: 8px; border-top: 1px solid #dde5e8;
         font-size: 9.5px; color: #8b9aa2; }
"""


def esc(v):
    return html_mod.escape("" if v is None else str(v))


def page(title, eyebrow, subtitle, body_html, source_note):
    return (
        "<!doctype html><html lang='ar' dir='rtl'><head><meta charset='utf-8'>"
        "<title>%s</title><style>%s</style></head><body>"
        "<header><div class='eyebrow'>%s</div><h1>%s</h1>"
        "<div class='sub'>%s</div></header>%s"
        "<footer>%s</footer></body></html>"
        % (esc(title), CSS, esc(eyebrow), esc(title), esc(subtitle), body_html,
           esc(source_note))
    )


def kpi(k, v, n=""):
    return ("<div class='kpi'><div class='k'>%s</div><div class='v'>%s</div>"
            "<div class='n'>%s</div></div>" % (esc(k), esc(v), esc(n)))


def render_pdf(html_text, out_path):
    if not CHROME:
        raise RuntimeError("google-chrome not found on this host")
    tmp = tempfile.mkdtemp(prefix="acms-rep-")
    try:
        src = os.path.join(tmp, "report.html")
        with open(src, "w", encoding="utf-8") as fh:
            fh.write(html_text)
        # A private profile dir keeps concurrent renders from fighting over one.
        cmd = [CHROME, "--headless", "--disable-gpu", "--no-sandbox",
               "--no-pdf-header-footer", "--virtual-time-budget=8000",
               "--user-data-dir=" + os.path.join(tmp, "profile"),
               "--print-to-pdf=" + out_path, "file://" + src]
        p = subprocess.run(cmd, capture_output=True, timeout=180)
        if not os.path.exists(out_path):
            raise RuntimeError("chrome produced no file: %s"
                               % p.stderr.decode()[-300:])
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ------------------------------------------------------------------- reports

def html_exec_summary(d, narrative):
    metrics = {m["code"]: m for m in d["metrics"]}

    def mval(code):
        m = metrics.get(code)
        if not m:
            return "—", ""
        if m.get("value") is None:
            return "—", "غير متاح: %s" % (m.get("unavailableReason") or "")
        unit = m.get("unit")
        v = m["value"]
        txt = money(v) if unit == "CURRENCY" else (
            pct(v) if unit == "PERCENT" else "{:,.2f}".format(num(v)).rstrip("0").rstrip("."))
        return txt, "على %s سجل" % m.get("basis", 0)

    kpis = "".join(kpi(METRIC_LABELS.get(c, c), *mval(c)) for c in
                   ("WEIGHTED_PIPELINE", "WIN_RATE", "GROSS_MARGIN",
                    "FORECAST_ACCURACY", "OPEN_APPROVALS", "APPROVAL_WAIT"))

    stage_rows = "".join(
        "<tr><td>%s</td><td class='n'>%d</td><td class='n'>%s</td></tr>"
        % (esc(STAGE_LABELS.get(s, s)), e["count"], money(e["value"]))
        for s, e in d["byStage"]) or "<tr><td colspan='3' class='empty'>لا توجد فرص مفتوحة</td></tr>"

    opp_rows = "".join(
        "<tr><td class='n'>%s</td><td>%s</td><td>%s</td><td class='n'>%s</td>"
        "<td class='n'>%s%%</td><td class='n'>%s</td></tr>"
        % (esc(o.get("code")), esc(o.get("name")),
           esc(STAGE_LABELS.get(o.get("stage"), o.get("stage"))),
           money(o.get("estimatedValue"), o.get("currency") or "USD"),
           "{:,.0f}".format(num(o.get("probability"))),
           esc((o.get("expectedCloseDate") or "—")[:10]))
        for o in d["active"]) or "<tr><td colspan='6' class='empty'>لا توجد فرص مفتوحة</td></tr>"

    dl_rows = "".join(
        "<tr><td class='n'>%s</td><td>%s</td><td class='n'>%s</td></tr>"
        % (esc(b.get("tenderNumber") or b.get("id", "")[:8]),
           esc((b.get("opportunity") or {}).get("name") or "—"),
           esc((b.get("submissionDeadline") or "—")[:10]))
        for b in d["deadlines"]) or (
        "<tr><td colspan='3' class='empty'>لا توجد مناقصات تُغلق خلال %d يومًا</td></tr>"
        % d["days"])

    narrative_html = ("<h2>Commentary</h2><div class='narrative'>%s</div>"
                      % esc(narrative).replace("\n", "<br>")) if narrative else ""

    body = (
        "<h2>Key figures</h2><div class='kpis'>%s</div>"
        "%s"
        "<h2>Pipeline by stage</h2><table><thead><tr><th>المرحلة</th><th>عدد</th>"
        "<th>القيمة التقديرية</th></tr></thead><tbody>%s"
        "<tr class='total'><td>الإجمالي المفتوح</td><td class='n'>%d</td>"
        "<td class='n'>%s</td></tr></tbody></table>"
        "<h2>Open opportunities</h2><table><thead><tr><th>الكود</th><th>الاسم</th>"
        "<th>المرحلة</th><th>القيمة</th><th>الاحتمال</th><th>الإغلاق المتوقع</th>"
        "</tr></thead><tbody>%s</tbody></table>"
        "<h2>Bid deadlines · %d days</h2><table><thead><tr><th>المناقصة</th>"
        "<th>الفرصة</th><th>آخر موعد</th></tr></thead><tbody>%s</tbody></table>"
        "<h2>Closed</h2><table><thead><tr><th>الحالة</th><th>عدد</th><th>القيمة</th>"
        "</tr></thead><tbody>"
        "<tr><td>مكسوبة</td><td class='n'>%d</td><td class='n'>%s</td></tr>"
        "<tr><td>مخسورة</td><td class='n'>%d</td><td class='n'>%s</td></tr>"
        "</tbody></table>"
        % (kpis, narrative_html, stage_rows, d["openCount"], money(d["openValue"]),
           opp_rows, d["days"], dl_rows,
           d["wonCount"], money(d["wonValue"]), d["lostCount"], money(d["lostValue"]))
    )
    return page("الملخّص التنفيذي — ACMS", "ACMS · EXECUTIVE SUMMARY",
                "حتى %s · مولَّد آليًا من بيانات النظام" % d["asOf"], body,
                "كل الأرقام مأخوذة من واجهة ACMS ومحسوبة في خدمة التقارير — لا تقديرات.")


def html_cost_sheet(d, narrative):
    o = d["opportunity"]
    cur = d["currency"]
    rows = "".join(
        "<tr><td>%s</td><td>%s</td><td class='n'>%s</td><td class='n'>%s</td>"
        "<td class='n'>%s</td><td class='n'>%s</td><td class='n'>%s</td></tr>"
        % (esc(r["package"]), esc(r["description"]), esc(r["unit"]),
           "{:,.2f}".format(r["quantity"]).rstrip("0").rstrip("."),
           "{:,.2f}".format(r["sellingRate"]),
           "{:,.2f}".format(r["linePrice"]), "{:,.2f}".format(r["lineCost"]))
        for r in d["rows"]) or "<tr><td colspan='7' class='empty'>لا توجد بنود تكلفة مسجَّلة</td></tr>"

    narrative_html = ("<h2>Commentary</h2><div class='narrative'>%s</div>"
                      % esc(narrative).replace("\n", "<br>")) if narrative else ""

    body = (
        "<h2>Key figures</h2><div class='kpis'>%s%s%s</div>%s"
        "<h2>Bill of quantities</h2><table><thead><tr><th>الحزمة</th><th>البند</th>"
        "<th>الوحدة</th><th>الكمية</th><th>سعر البيع</th><th>إجمالي البيع</th>"
        "<th>إجمالي التكلفة</th></tr></thead><tbody>%s"
        "<tr class='total'><td colspan='5'>الإجمالي</td><td class='n'>%s</td>"
        "<td class='n'>%s</td></tr></tbody></table>"
        % (kpi("إجمالي سعر البيع", money(d["grandPrice"], cur)),
           kpi("إجمالي التكلفة", money(d["grandCost"], cur)),
           kpi("الهامش", pct(d["margin"]), "محسوب على سعر البيع"),
           narrative_html, rows,
           "{:,.2f}".format(d["grandPrice"]), "{:,.2f}".format(d["grandCost"]))
    )
    return page("ورقة التكلفة — %s" % (o.get("name") or "—"), "ACMS · COST SHEET",
                "%s · سيناريو: %s · حتى %s"
                % (o.get("code") or "—", d["scenario"].get("name") or "—", d["asOf"]),
                body,
                "الهامش يُحسب على سعر البيع لا على التكلفة، وفق قاعدة ACMS نفسها.")


# ------------------------------------------------------------------ PPTX

def build_pptx(title, subtitle, slides, out_path):
    """slides: list of (heading, [(label, value)] or [str] bullets)."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    INK = RGBColor(0x10, 0x18, 0x1C)
    SIGNAL = RGBColor(0x0D, 0x6F, 0x7C)
    MUTED = RGBColor(0x6F, 0x80, 0x89)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def rtl(paragraph):
        # python-pptx has no RTL switch; the attribute has to go on the XML.
        paragraph.alignment = PP_ALIGN.RIGHT
        paragraph._pPr.set("rtl", "1")

    def textbox(slide, left, top, width, height):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    # --- title slide
    s = prs.slides.add_slide(blank)
    tf = textbox(s, 0.7, 2.4, 12.0, 2.4)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = INK
    r.font.name = "Segoe UI"
    rtl(p)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(16); r2.font.color.rgb = MUTED; r2.font.name = "Segoe UI"
    rtl(p2)

    # --- content slides
    for heading, items in slides:
        s = prs.slides.add_slide(blank)
        htf = textbox(s, 0.7, 0.5, 12.0, 0.9)
        hp = htf.paragraphs[0]
        hr = hp.add_run(); hr.text = heading
        hr.font.size = Pt(26); hr.font.bold = True; hr.font.color.rgb = SIGNAL
        hr.font.name = "Segoe UI"
        rtl(hp)

        btf = textbox(s, 0.7, 1.6, 12.0, 5.2)
        first = True
        for it in items:
            p = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            if isinstance(it, tuple):
                label, value = it
                r = p.add_run(); r.text = "%s   " % label
                r.font.size = Pt(16); r.font.color.rgb = INK; r.font.name = "Segoe UI"
                rv = p.add_run(); rv.text = str(value)
                rv.font.size = Pt(16); rv.font.bold = True
                rv.font.color.rgb = SIGNAL; rv.font.name = "Consolas"
            else:
                r = p.add_run(); r.text = "•  %s" % it
                r.font.size = Pt(15); r.font.color.rgb = INK; r.font.name = "Segoe UI"
            p.space_after = Pt(9)
            rtl(p)

    prs.save(out_path)


def pptx_exec_summary(d, narrative, out_path):
    metrics = {m["code"]: m for m in d["metrics"]}

    def mtxt(code):
        m = metrics.get(code)
        if not m or m.get("value") is None:
            return "غير متاح"
        return (money(m["value"]) if m.get("unit") == "CURRENCY"
                else pct(m["value"]) if m.get("unit") == "PERCENT"
                else "{:,.2f}".format(num(m["value"])).rstrip("0").rstrip("."))

    slides = [
        ("مؤشرات الأداء", [(METRIC_LABELS[c], mtxt(c)) for c in METRIC_LABELS]),
        ("الـPipeline حسب المرحلة",
         [("%s — %d فرصة" % (STAGE_LABELS.get(s, s), e["count"]), money(e["value"]))
          for s, e in d["byStage"]] or ["لا توجد فرص مفتوحة"]),
        ("أكبر الفرص المفتوحة",
         [("%s — %s" % (o.get("code"), o.get("name")),
           money(o.get("estimatedValue"), o.get("currency") or "USD"))
          for o in d["active"][:8]] or ["لا توجد فرص مفتوحة"]),
        ("مواعيد المناقصات · %d يومًا" % d["days"],
         ["%s — %s" % ((b.get("tenderNumber") or "—"),
                       (b.get("submissionDeadline") or "—")[:10])
          for b in d["deadlines"]] or ["لا توجد مناقصات تُغلق في هذه المدة"]),
        ("الحصيلة",
         [("مفتوحة", "%d · %s" % (d["openCount"], money(d["openValue"]))),
          ("مكسوبة", "%d · %s" % (d["wonCount"], money(d["wonValue"]))),
          ("مخسورة", "%d · %s" % (d["lostCount"], money(d["lostValue"]))),
          ("موافقات معلّقة", str(len(d["approvals"])))]),
    ]
    if narrative:
        slides.append(("قراءة تنفيذية",
                       [ln for ln in narrative.split("\n") if ln.strip()][:8]))

    build_pptx("الملخّص التنفيذي — ACMS", "حتى %s · مولَّد آليًا" % d["asOf"],
               slides, out_path)


def pptx_cost_sheet(d, narrative, out_path):
    o = d["opportunity"]
    cur = d["currency"]
    by_pkg = {}
    for r in d["rows"]:
        e = by_pkg.setdefault(r["package"], {"price": 0.0, "cost": 0.0, "n": 0})
        e["price"] += r["linePrice"]
        e["cost"] += r["lineCost"]
        e["n"] += 1

    slides = [
        ("الأرقام الرئيسية",
         [("إجمالي سعر البيع", money(d["grandPrice"], cur)),
          ("إجمالي التكلفة", money(d["grandCost"], cur)),
          ("الهامش (على سعر البيع)", pct(d["margin"]))]),
        ("الحزم",
         [("%s — %d بند" % (name, e["n"]), money(e["price"], cur))
          for name, e in by_pkg.items()] or ["لا توجد حزم"]),
        ("أكبر البنود",
         [("%s" % r["description"], money(r["linePrice"], cur))
          for r in sorted(d["rows"], key=lambda r: -r["linePrice"])[:8]]
         or ["لا توجد بنود"]),
    ]
    if narrative:
        slides.append(("ملاحظات", [ln for ln in narrative.split("\n") if ln.strip()][:8]))

    build_pptx("ورقة التكلفة — %s" % (o.get("name") or "—"),
               "%s · %s" % (o.get("code") or "—", d["asOf"]), slides, out_path)


REPORTS = ("exec_summary", "cost_sheet")
FORMATS = ("pdf", "pptx")


def generate(payload):
    report = payload.get("report", "exec_summary")
    fmt = payload.get("format", "pdf")
    narrative = (payload.get("narrative") or "").strip()

    if report not in REPORTS:
        raise ValueError("report must be one of %s" % ", ".join(REPORTS))
    if fmt not in FORMATS:
        raise ValueError("format must be one of %s" % ", ".join(FORMATS))

    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    name = "%s-%s-%s.%s" % (report.replace("_", "-"), stamp, uuid.uuid4().hex[:6], fmt)
    out = os.path.join(FILES, name)

    if report == "exec_summary":
        d = gather_summary(days=payload.get("days", 60))
        if fmt == "pdf":
            render_pdf(html_exec_summary(d, narrative), out)
        else:
            pptx_exec_summary(d, narrative, out)
        facts = {
            "openCount": d["openCount"], "openValue": round(d["openValue"], 2),
            "weightedPipeline": round(d["weighted"], 2),
            "wonCount": d["wonCount"], "lostCount": d["lostCount"],
            "deadlines": len(d["deadlines"]), "pendingApprovals": len(d["approvals"]),
        }
    else:
        opp_id = payload.get("opportunityId")
        if not opp_id:
            raise ValueError("cost_sheet needs opportunityId")
        d = gather_cost_sheet(opp_id)
        if fmt == "pdf":
            render_pdf(html_cost_sheet(d, narrative), out)
        else:
            pptx_cost_sheet(d, narrative, out)
        facts = {
            "opportunity": (d["opportunity"] or {}).get("code"),
            "lineItems": len(d["rows"]),
            "totalPrice": round(d["grandPrice"], 2),
            "totalCost": round(d["grandCost"], 2),
            "marginPercent": None if d["margin"] is None else round(d["margin"], 2),
        }

    return {
        "ok": True,
        "report": report,
        "format": fmt,
        "fileName": name,
        "url": "%s/files/%s" % (PUBLIC_BASE.rstrip("/"), name),
        "sizeBytes": os.path.getsize(out),
        "facts": facts,
    }


# --------------------------------------------------------------------- HTTP

# ------------------------------------------------------- read projections
#
# n8n hands an httpRequestTool's HTTP response straight to the model — there is
# no hook in between where a Code node could trim it. So a 15-record
# opportunity list reached the LLM as 17,770 characters of 38-field records
# (~4,400 tokens), and two such calls put a specialist past Groq's 8,000 TPM
# ceiling. The projection therefore has to happen at the source, which is here.
# These endpoints call the same ACMS API with the caller's own bearer token, so
# who may see what is unchanged — only how much of it crosses into the prompt.

Q_LIMIT_DEFAULT = 20
Q_LIMIT_MAX = 100

# The ceiling a single tool response may occupy in the prompt. ~6,000 characters
# is roughly 1,500 tokens, so even four reads in one agent turn stay well inside
# Groq's 8,000 TPM budget with room for the system prompt and the answer.
Q_MAX_CHARS = 6000

# Flattened so a record stays one level deep: a nested object costs braces,
# quotes and its own repeated id for one useful string.
_FLAT = {
    "accountName": ("account", "legalName"),
    "accountCode": ("account", "code"),
    "ownerName": ("owner", "fullNameEn"),
    "userName": ("user", "fullNameEn"),
    "opportunityName": ("opportunity", "name"),
    "opportunityCode": ("opportunity", "code"),
}

Q_SPECS = {
    "opportunities": {
        "path": "/opportunities",
        # Only the params ACMS itself accepts. It rejects anything unknown with
        # a 400 ("property limit should not exist"), so the rest is applied here.
        "passthrough": ("status", "stage", "country", "search"),
        "local": ("industry",),
        "views": {
            "list": ("id", "code", "name", "stage", "status", "estimatedValue",
                     "currency", "expectedCloseDate", "accountName"),
            "forecast": ("id", "code", "name", "stage", "status", "estimatedValue",
                         "currency", "expectedCloseDate", "probability",
                         "forecastCategory", "health", "accountName"),
            "detail": ("id", "code", "name", "stage", "status", "estimatedValue",
                       "proposedPrice", "estimatedCost", "marginPercent", "currency",
                       "expectedCloseDate", "probability", "forecastCategory",
                       "health", "country", "industry", "source", "nextStep",
                       "accountName", "ownerName"),
        },
    },
    "activities": {
        "path": "/activities",
        "passthrough": ("type", "opportunityId"),
        "local": ("completed",),
        "views": {
            "list": ("id", "type", "subject", "dueAt", "completedAt",
                     "opportunityCode", "opportunityName", "accountName", "userName"),
            "detail": ("id", "type", "subject", "body", "dueAt", "completedAt",
                       "createdAt", "opportunityId", "opportunityCode",
                       "opportunityName", "accountName", "userName"),
        },
    },
    "accounts": {
        "path": "/accounts",
        "passthrough": ("search", "country", "type"),
        "local": (),
        "views": {
            "list": ("id", "code", "legalName", "tradeName", "country", "type",
                     "industry", "creditStatus"),
            "detail": ("id", "code", "legalName", "tradeName", "country", "type",
                       "industry", "creditStatus", "city", "website",
                       "paymentTermDays", "ownerName"),
        },
    },
}


# ------------------------------------------------------- session business context
#
# What the conversation is currently *about* — the opportunity or account in
# play — kept as a field rather than as a sentence in the transcript. Its life
# is the working session, so the window is hours rather than the ten minutes a
# proposal gets.

CONTEXT_TTL = 6 * 60 * 60
_context = {}
_context_lock = threading.Lock()


def _prune_context():
    now = time.time()
    for k in [k for k, v in _context.items() if now - v.get("ts", 0) > CONTEXT_TTL]:
        _context.pop(k, None)


def context_put(session_id, patch):
    """Merge, never replace: a turn that names only an account must not erase
    the opportunity the previous turn established."""
    with _context_lock:
        _prune_context()
        cur = _context.get(session_id) or {}
        for k, v in (patch or {}).items():
            if v not in (None, "", [], {}):
                cur[k] = v
        cur["ts"] = time.time()
        cur["updatedAt"] = datetime.now(timezone.utc).isoformat(timespec="seconds")
        _context[session_id] = cur
        return {k: v for k, v in cur.items() if k != "ts"}


def context_get(session_id):
    with _context_lock:
        _prune_context()
        cur = _context.get(session_id)
        return {k: v for k, v in cur.items() if k != "ts"} if cur else None


# --------------------------------------------------------------- answer cache
#
# Measured before it was built, on genuine chat-widget traffic only — developer
# sessions and the regression suite repeat questions by construction and would
# have proved the cache works against my own test loop. Over 138 real turns the
# repeat rate for read questions was 21% at a 15-minute window, 27% at an hour,
# and 34% at a day. An hour takes most of the benefit; the seven further points
# up to a day cost a whole day of possible staleness, so an hour it is.
ANSWER_CACHE_TTL = int(os.environ.get("ACMS_CACHE_TTL", "3600"))

# TTL alone would still serve an answer whose data changed a minute after it was
# stored. Every entry therefore carries a fingerprint of the data it was derived
# from, and a hit requires that fingerprint to still match. Recomputing it costs
# three cheap list calls against a model call of several seconds, and the result
# is memoised briefly so a burst of questions does not repeat them.
DATA_VERSION_TTL = int(os.environ.get("ACMS_CACHE_DV_TTL", "30"))
DATA_VERSION_SOURCES = ("/opportunities", "/accounts", "/activities")

_answer_cache = {}
_answer_cache_lock = threading.Lock()
_dv_cache = {"v": None, "ts": 0.0}
_dv_lock = threading.Lock()

CACHE_STATS = {"lookups": 0, "hits": 0, "similar": 0, "stale": 0, "expired": 0, "stores": 0}


def _cache_norm(q):
    """Same normalisation the gate uses, so 'الفرص' and 'الفُرص' are one key."""
    t = re.sub(r"[\u064b-\u0652\u0640]", "", str(q or ""))
    t = re.sub(r"[\u0625\u0623\u0622\u0627]", "\u0627", t)
    t = t.replace("\u0649", "\u064a").replace("\u0629", "\u0647")
    t = re.sub(r"[^\w\s]", " ", t, flags=re.UNICODE)
    return " ".join(t.split()).strip().lower()


def _cache_key(user, question):
    # The user is part of the key, not a detail. Two people can ask the same
    # sentence and be entitled to different rows, so an answer is only ever
    # replayed to the identity it was produced for. While identity mode is
    # "optional" everyone unbound shares one bucket — which is correct, because
    # they genuinely share one ACMS identity today; the buckets separate by
    # themselves the moment the mode becomes "required".
    raw = (str(user or "svc") + "\n" + _cache_norm(question)).encode("utf-8")
    return hashlib.sha256(raw).hexdigest()


def _data_version():
    """A fingerprint that changes whenever the answerable data changes."""
    now = time.time()
    with _dv_lock:
        if _dv_cache["v"] is not None and now - _dv_cache["ts"] < DATA_VERSION_TTL:
            return _dv_cache["v"]
    parts = []
    for path in DATA_VERSION_SOURCES:
        try:
            r = acms.get(path)
        except Exception:
            # Unknown beats wrong: if the marker cannot be computed, no entry
            # may be served from cache this turn.
            return None
        items = (r or {}).get("items", r if isinstance(r, list) else []) or []
        newest = ""
        for it in items:
            u = str(it.get("updatedAt") or it.get("createdAt") or "")
            if u > newest:
                newest = u
        parts.append("%s:%d:%s" % (path, len(items), newest))
    v = hashlib.sha256("|".join(parts).encode("utf-8")).hexdigest()[:16]
    with _dv_lock:
        _dv_cache["v"] = v
        _dv_cache["ts"] = now
    return v


# Words that carry no topic. Dropping them is what lets "اعرض الفرص" and
# "ممكن الفرص لو سمحت" be one question; keeping them would make politeness a
# cache miss.
CACHE_STOPWORDS = {
    "ايه", "ما", "ماهي", "هو", "هي", "من", "في", "علي", "عن", "الي", "مع",
    "دلوقتي", "حاليا", "الان", "ممكن", "لو", "سمحت", "عايز", "عاوز", "اريد",
    # Verb forms of the same request, all of which mean "show me". Listing one
    # and not its siblings made "ممكن تعرض" score 0.80 against "اعرض" and miss.
    "اعرض", "تعرض", "عرض", "اعرضلي", "وريني", "قولي", "تقولي", "قوللي",
    "هات", "اعطني", "تديني", "كل", "اللي", "التي", "الذي",
    "و", "يا", "هل", "فيه", "في", "بتاع", "بتاعت", "شوف",
}

# 0.82 is deliberately high. The measurement showed 0.6 buys one extra point
# over 0.8 while letting genuinely different questions match, and a wrong
# cached answer is far worse than a missed one: it is confidently wrong, it
# cites real record codes, and nothing downstream would flag it.
CACHE_SIM_THRESHOLD = float(os.environ.get("ACMS_CACHE_SIM", "0.82"))


def _cache_bag(normalised):
    return frozenset(w for w in normalised.split()
                     if w and w not in CACHE_STOPWORDS)


def _cache_pins(normalised):
    """Tokens that must match exactly, whatever the similarity score says.

    Two questions differing only by a number or a code are different questions:
    "خلال 60 يوم" against "خلال 30 يوم", or one opportunity code against
    another. Jaccard treats a single differing token as noise at these lengths,
    which is exactly the case where it must not.
    """
    pins = set()
    for w in normalised.split():
        if any(ch.isdigit() for ch in w):
            pins.add(w)
        elif w.isascii() and len(w) >= 3:      # FTTH, OPP, ACC, backbone …
            pins.add(w)
    return frozenset(pins)


def _similar(a, b):
    """Jaccard over topic words, with the pinned tokens required to agree."""
    if a["pins"] != b["pins"]:
        return 0.0
    x, y = a["bag"], b["bag"]
    if not x or not y:
        return 0.0
    inter = len(x & y)
    union = len(x | y)
    return inter / union if union else 0.0


def _prune_answer_cache():
    now = time.time()
    for k in [k for k, v in _answer_cache.items()
              if now - v.get("ts", 0) > ANSWER_CACHE_TTL]:
        _answer_cache.pop(k, None)


def cache_lookup(user, question):
    CACHE_STATS["lookups"] += 1
    dv = _data_version()
    if dv is None:
        return {"hit": False, "reason": "no-data-version"}
    key = _cache_key(user, question)
    with _answer_cache_lock:
        # Read the entry before pruning, not after. Pruning first deletes an
        # expired entry and the lookup then reports "miss", so the expired
        # counter could never move and the stats would misattribute every
        # timeout as a question nobody had asked before.
        e = _answer_cache.get(key)
        _prune_answer_cache()
        matched = "exact"
        if not e:
            # No identical phrasing. Look for the same question worded
            # differently — same user only, since the user is part of the key
            # for a reason and similarity must not become a way around it.
            n = _cache_norm(question)
            probe = {"bag": _cache_bag(n), "pins": _cache_pins(n)}
            best, best_score = None, 0.0
            uid = str(user or "svc")
            for cand in _answer_cache.values():
                if cand.get("user") != uid:
                    continue
                sc = _similar(probe, cand)
                if sc > best_score:
                    best, best_score = cand, sc
            if best is not None and best_score >= CACHE_SIM_THRESHOLD:
                e, matched = best, "similar:%.2f" % best_score
            else:
                return {"hit": False, "reason": "miss"}
        if time.time() - e["ts"] > ANSWER_CACHE_TTL:
            _answer_cache.pop(_cache_key(e.get("user"), e.get("q", "")), None)
            _answer_cache.pop(key, None)
            CACHE_STATS["expired"] += 1
            return {"hit": False, "reason": "expired"}
        if e.get("dv") != dv:
            # The data moved. Drop it rather than keep serving a stale answer.
            _answer_cache.pop(_cache_key(e.get("user"), e.get("q", "")), None)
            _answer_cache.pop(key, None)
            CACHE_STATS["stale"] += 1
            return {"hit": False, "reason": "data-changed"}
        CACHE_STATS["hits"] += 1
        if matched != "exact":
            CACHE_STATS["similar"] = CACHE_STATS.get("similar", 0) + 1
        return {"hit": True, "answer": e["answer"], "matched": matched,
                "ageMin": int((time.time() - e["ts"]) / 60)}


def cache_store(user, question, answer):
    dv = _data_version()
    if dv is None:
        return {"ok": False, "reason": "no-data-version"}
    text = str(answer or "")
    if len(text) < 20:
        return {"ok": False, "reason": "too-short"}
    with _answer_cache_lock:
        _prune_answer_cache()
        n = _cache_norm(question)
        _answer_cache[_cache_key(user, question)] = {
            "answer": text, "ts": time.time(), "dv": dv,
            "user": str(user or "svc"), "q": n,
            "bag": _cache_bag(n), "pins": _cache_pins(n)}
        CACHE_STATS["stores"] += 1
    return {"ok": True}


# ------------------------------------------------------------- telemetry read
#
# The store is written by ~/acms-telemetry/collect.py, which harvests n8n's own
# execution records. Nothing is instrumented inside the workflow: n8n already
# keeps token counts and timings, and harvesting afterwards also captures the
# runs that failed — the ones an in-workflow emitter would miss.

TELEMETRY = os.environ.get(
    "ACMS_TELEMETRY", os.path.expanduser("~/acms-telemetry/telemetry.jsonl"))

# Unset by default and reported as null rather than guessed.
PRICE_IN = float(os.environ.get("ACMS_PRICE_IN_PER_MTOK", "0") or 0)
PRICE_OUT = float(os.environ.get("ACMS_PRICE_OUT_PER_MTOK", "0") or 0)


def _pct(values, p):
    if not values:
        return None
    s = sorted(values)
    return s[min(len(s) - 1, int(len(s) * p))]


def telemetry_summary(days):
    if not os.path.exists(TELEMETRY):
        return {"error": "no telemetry store at %s" % TELEMETRY}

    cutoff = datetime.now(timezone.utc) - timedelta(days=days)
    recs = []
    with open(TELEMETRY, encoding="utf-8") as fh:
        for line in fh:
            line = line.strip()
            if not line:
                continue
            try:
                r = json.loads(line)
            except ValueError:
                continue
            at = _iso((r.get("at") or "").replace(" ", "T") + "+00:00")
            if at and at < cutoff:
                continue
            recs.append(r)

    if not recs:
        return {"window_days": days, "executions": 0}

    tok = [r.get("totalTokens") or 0 for r in recs]
    ms = [r["ms"] for r in recs if isinstance(r.get("ms"), int)]
    routes, intents, errors, disp = {}, {}, {}, {}
    limited = 0
    prompt_sum = out_sum = 0
    for r in recs:
        routes[r.get("route") or "?"] = routes.get(r.get("route") or "?", 0) + 1
        intents[r.get("intent") or "?"] = intents.get(r.get("intent") or "?", 0) + 1
        if r.get("limited"):
            limited += 1
        if r.get("errorClass"):
            errors[r["errorClass"]] = errors.get(r["errorClass"], 0) + 1
        d = str(r.get("dispatchCount", 0))
        disp[d] = disp.get(d, 0) + 1
        prompt_sum += r.get("promptTokens") or 0
        out_sum += r.get("completionTokens") or 0

    cost = None
    if PRICE_IN or PRICE_OUT:
        cost = round(prompt_sum / 1e6 * PRICE_IN + out_sum / 1e6 * PRICE_OUT, 4)

    # A turn that never reached a model is the fast path doing its job.
    no_model = sum(1 for r in recs if not r.get("modelCalls"))

    return {
        "window_days": days,
        "executions": len(recs),
        "tokens": {
            "prompt": prompt_sum,
            "completion": out_sum,
            "total": prompt_sum + out_sum,
            "per_run_p50": _pct(tok, 0.5),
            "per_run_p95": _pct(tok, 0.95),
            "per_run_max": max(tok) if tok else 0,
        },
        "latency_ms": {"p50": _pct(ms, 0.5), "p95": _pct(ms, 0.95),
                       "max": max(ms) if ms else None},
        "routes": routes,
        "intents": intents,
        "dispatch_histogram": disp,
        "no_model_turns": no_model,
        "no_model_share": round(no_model / len(recs), 3),
        "rate_limited": limited,
        "rate_limited_share": round(limited / len(recs), 3),
        "errors": errors,
        # Null unless a price is configured — tokens are measured, a price is not.
        "cost": cost,
        "cost_note": None if cost is not None else
        "اضبط ACMS_PRICE_IN_PER_MTOK و ACMS_PRICE_OUT_PER_MTOK لحساب التكلفة",
    }


# ------------------------------------------------------------------ svc auth
#
# One shared secret for the endpoints only n8n calls. Read from a file so it
# never sits in the process list or in this source.

def _load_svc_token():
    path = os.environ.get("ACMS_SVC_TOKEN_FILE",
                          os.path.join(HERE, ".svc-token"))
    try:
        with open(path, encoding="utf-8") as fh:
            return fh.read().strip() or None
    except OSError:
        return None


SVC_TOKEN = _load_svc_token()

# Reachable without the shared secret. `/files` serves a generated report to the
# user's browser and validates the filename already; `/health` must stay probe-able.
OPEN_PREFIXES = ("/health", "/files/", "/login", "/identity/start")


def _svc_ok(handler):
    """True when the caller proved it is the workflow, or the path is open."""
    if SVC_TOKEN is None:
        return True                      # not configured — fail open, and say so
    p = handler.path
    if p == "/" or any(p.startswith(x) for x in OPEN_PREFIXES):
        return True
    if p.startswith("/q/"):
        return True                      # guarded by the caller's ACMS bearer
    return handler.headers.get("X-ACMS-Svc") == SVC_TOKEN


# ------------------------------------------------------------------ identity
#
# Every chat user shared one service account, so ACMS saw a single CEO-scoped
# caller no matter who was typing — permissions were the account's, not the
# person's, and the audit trail named the robot. This binds a chat session to a
# real ACMS user.
#
# ACMS_IDENTITY_MODE:
#   optional  bound sessions use their own token; unbound ones fall back to the
#             service account, so nothing breaks while people are still signing
#             in. Every resolution is recorded either way.
#   required  an unbound session is refused and told to sign in.
IDENTITY_MODE = os.environ.get("ACMS_IDENTITY_MODE", "optional")

IDENTITY_TTL = 12 * 60 * 60      # a working day
BIND_TTL = 120                   # the code is a handoff, not a credential

_identity = {}                   # sessionId -> {token, email, sub, roles, ts}
_bindcodes = {}                  # code      -> {token, email, sub, roles, ts}
_identity_lock = threading.Lock()


def _prune_identity():
    now = time.time()
    for k in [k for k, v in _identity.items() if now - v["ts"] > IDENTITY_TTL]:
        _identity.pop(k, None)
    for k in [k for k, v in _bindcodes.items() if now - v["ts"] > BIND_TTL]:
        _bindcodes.pop(k, None)


def _claims(token):
    """Read sub/email/roles out of the ACMS JWT. Signature is ACMS's business —
    this only needs to know who the token says it is, for display and audit."""
    try:
        part = token.split(".")[1]
        part += "=" * (-len(part) % 4)
        return json.loads(base64.urlsafe_b64decode(part))
    except Exception:
        return {}


def identity_start(email, password):
    """Exchange credentials for an ACMS token and hand back a binding code.

    The password exists only inside this call. It is not stored, not returned,
    and not written to the log — the access log records the path alone.
    """
    if not email or not password:
        return 400, {"error": "email and password are required"}
    body = json.dumps({"email": email, "password": password}).encode()
    req = urllib.request.Request(API + "/auth/login", data=body, method="POST")
    req.add_header("Content-Type", "application/json")
    try:
        with urllib.request.urlopen(req, timeout=30) as r:
            tok = json.load(r).get("accessToken")
    except urllib.error.HTTPError:
        # Deliberately not distinguishing "no such user" from "wrong password".
        return 401, {"error": "بيانات الدخول غير صحيحة"}
    except Exception:
        return 502, {"error": "تعذّر الوصول إلى ACMS"}
    if not tok:
        return 502, {"error": "ACMS returned no token"}

    c = _claims(tok)
    code = "".join(random.choice("ABCDEFGHJKLMNPQRSTUVWXYZ23456789") for _ in range(6))
    with _identity_lock:
        _prune_identity()
        _bindcodes[code] = {"token": tok, "email": c.get("email") or email,
                            "sub": c.get("sub"), "roles": c.get("roles"),
                            "ts": time.time()}
    return 200, {"ok": True, "code": code, "expires_in": BIND_TTL,
                 "email": c.get("email") or email}


def identity_bind(code, session_id):
    with _identity_lock:
        _prune_identity()
        rec = _bindcodes.pop(str(code).strip().upper(), None)
        if not rec:
            return 404, {"ok": False, "error": "رمز غير صالح أو منتهٍ"}
        rec = dict(rec, ts=time.time())
        _identity[session_id] = rec
    return 200, {"ok": True, "email": rec["email"], "roles": rec.get("roles")}


def identity_get(session_id):
    with _identity_lock:
        _prune_identity()
        return _identity.get(session_id)


def identity_drop(session_id):
    with _identity_lock:
        return _identity.pop(session_id, None) is not None


def identity_resolve(session_id):
    """What the workflow calls at the start of every turn."""
    rec = identity_get(session_id)
    if rec:
        return 200, {"identified": True, "accessToken": rec["token"],
                     "email": rec["email"], "sub": rec.get("sub"),
                     "roles": rec.get("roles"), "mode": IDENTITY_MODE}
    if IDENTITY_MODE == "required":
        return 200, {"identified": False, "accessToken": None,
                     "mode": "required",
                     "message": "سجّل الدخول أولًا: افتح صفحة الدخول ثم اكتب هنا "
                                "/login ورمز الربط."}
    # optional: the service account carries the turn, and the answer says so.
    return 200, {"identified": False, "accessToken": acms.token(),
                 "email": EMAIL, "sub": None, "roles": None,
                 "mode": "optional",
                 "message": "جلسة غير مُوثّقة — تعمل بحساب الخدمة."}


LOGIN_PAGE = """<!doctype html><html lang="ar" dir="rtl"><head>
<meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>ACMS — تسجيل الدخول للمساعد</title>
<style>
 body{font-family:"Segoe UI",Tahoma,sans-serif;background:#FAF8F4;color:#1B1E22;
      display:flex;min-height:100vh;align-items:center;justify-content:center;margin:0}
 .card{background:#fff;border:1px solid #E2DED6;padding:30px 32px;width:min(94vw,420px)}
 h1{font-size:20px;margin:0 0 4px}
 p.sub{color:#6E7176;font-size:14px;margin:0 0 22px}
 label{display:block;font-size:13px;color:#3A4048;margin:14px 0 5px}
 input{width:100%;padding:10px 12px;border:1px solid #E2DED6;font-size:15px;
       font-family:inherit;box-sizing:border-box;background:#fff;color:#1B1E22}
 button{margin-top:20px;width:100%;padding:11px;border:0;background:#0072B2;color:#fff;
        font-size:15px;font-weight:600;cursor:pointer;font-family:inherit}
 button:disabled{opacity:.6;cursor:default}
 .out{margin-top:20px;padding:14px;border:1px solid #E2DED6;background:#F2EEE6;display:none}
 .code{font-family:ui-monospace,Consolas,monospace;font-size:30px;letter-spacing:.16em;
       font-weight:700;text-align:center;margin:8px 0;color:#0072B2}
 .err{color:#C0392B;font-size:14px;margin-top:14px;display:none}
 .note{color:#6E7176;font-size:12.5px;margin-top:18px;line-height:1.7}
</style></head><body>
<div class="card">
  <h1>تسجيل الدخول لمساعد ACMS</h1>
  <p class="sub">بحسابك أنت — ليصبح ما تراه وما تستطيع تنفيذه بصلاحياتك.</p>
  <form id="f" autocomplete="off">
    <label for="e">البريد</label>
    <input id="e" type="email" required autocomplete="username">
    <label for="p">كلمة المرور</label>
    <input id="p" type="password" required autocomplete="current-password">
    <button id="b" type="submit">احصل على رمز الربط</button>
  </form>
  <div class="err" id="err"></div>
  <div class="out" id="out">
    <div style="font-size:13px;color:#6E7176">اكتب في المحادثة:</div>
    <div class="code" id="code"></div>
    <div style="font-size:13px;color:#6E7176" id="who"></div>
  </div>
  <p class="note">
    كلمة المرور تُرسَل إلى خادم ACMS مباشرة ولا تُخزَّن ولا تُسجَّل ولا تمرّ بالمحادثة.
    الرمز صالح دقيقتين ويُستخدَم مرة واحدة.
  </p>
</div>
<script>
document.getElementById('f').addEventListener('submit', async function (ev) {
  ev.preventDefault();
  var b = document.getElementById('b'), err = document.getElementById('err');
  b.disabled = true; err.style.display = 'none';
  try {
    var r = await fetch('/identity/start', {
      method: 'POST', headers: {'Content-Type': 'application/json'},
      body: JSON.stringify({email: document.getElementById('e').value,
                            password: document.getElementById('p').value})});
    var d = await r.json();
    if (!r.ok) { throw new Error(d.error || 'تعذّر تسجيل الدخول'); }
    document.getElementById('code').textContent = '/login ' + d.code;
    document.getElementById('who').textContent = d.email;
    document.getElementById('out').style.display = 'block';
    document.getElementById('f').style.display = 'none';
  } catch (e) {
    err.textContent = e.message; err.style.display = 'block';
  }
  b.disabled = false;
});
</script></body></html>"""


# ------------------------------------------------------------ evidence ledger
#
# What each session has actually been shown. Written by /q, read by the output
# guard. Bounded and short-lived: it exists to check one answer, not to archive.

EVIDENCE_TTL = 30 * 60
_evidence = {}
_evidence_lock = threading.Lock()

# ACMS record codes look like OPP-2026-000013 / ACC-2026-000006.
CODE_RE = re.compile(r"\b[A-Z]{2,4}-\d{4}-\d{4,8}\b")
UUID_ANY = re.compile(
    r"\b[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}\b", re.I)


def _prune_evidence():
    now = time.time()
    for k in [k for k, v in _evidence.items() if now - v.get("ts", 0) > EVIDENCE_TTL]:
        _evidence.pop(k, None)


def evidence_record(session_id, resource, payload):
    """Note what was served. Codes and ids are what an answer can be checked against."""
    if not session_id:
        return
    blob = json.dumps(payload, ensure_ascii=False)
    with _evidence_lock:
        _prune_evidence()
        cur = _evidence.setdefault(session_id, {"codes": set(), "ids": set(),
                                                "sources": [], "ts": 0})
        cur["codes"].update(CODE_RE.findall(blob))
        cur["ids"].update(x.lower() for x in UUID_ANY.findall(blob))
        cur["sources"].append({
            "source": "ACMS", "endpoint": "/q/" + resource,
            "returned": payload.get("returned"), "total": payload.get("total"),
            "truncated": payload.get("truncated"),
            "facts": payload.get("facts"),
        })
        cur["sources"] = cur["sources"][-12:]
        cur["ts"] = time.time()


def evidence_get(session_id):
    with _evidence_lock:
        _prune_evidence()
        cur = _evidence.get(session_id)
        if not cur:
            return None
        return {"codes": sorted(cur["codes"]), "ids": sorted(cur["ids"]),
                "sources": cur["sources"]}


# ------------------------------------------------------------------ guardrails

# Instructions aimed at the model rather than at a reader. A tender PDF is the
# realistic carrier: it is long, nobody reads all of it, and it goes into the
# prompt whole.
INJECTION_PATTERNS = [
    (r"ignore (all|any|the) (previous|prior|above)", "ignore-previous"),
    (r"disregard .{0,20}(instructions|rules|prompt)", "disregard-rules"),
    (r"(you are|act as) (now )?(a |an )?(different|new) (assistant|ai|system)", "role-override"),
    (r"system\s*prompt", "system-prompt-probe"),
    (r"reveal .{0,25}(prompt|instructions|token|key|password)", "exfiltration"),
    (r"(execute|run|perform) .{0,25}without .{0,20}(confirmation|approval)", "bypass-approval"),
    (r"تجاهل .{0,20}(التعليمات|القواعد|السابق)", "ignore-previous-ar"),
    (r"نفّ?ذ .{0,25}(بدون|من غير) .{0,15}(تأكيد|موافقة)", "bypass-approval-ar"),
    (r"اكشف .{0,25}(التعليمات|المفتاح|كلمة|التوكن)", "exfiltration-ar"),
]

# Anything shaped like a credential must never reach an answer.
SECRET_PATTERNS = [
    # A real JWT header segment is ~15-30 chars after "eyJ"; requiring 20 let a
    # genuine token through the check.
    (r"eyJ[A-Za-z0-9_-]{10,}\.[A-Za-z0-9_-]{10,}", "jwt"),
    (r"\bBearer\s+[A-Za-z0-9._-]{20,}", "bearer"),
    (r"X-ACMS-Svc", "service-token-header"),
    (r"AcmsAgent#\S+", "service-password"),
    (r"\bsk-[A-Za-z0-9]{20,}", "api-key"),
]

CLAIMED_EXECUTION = re.compile(
    r"(تم\s+(الإنشاء|إنشاء|التحديث|تحديث|التنفيذ)|أنشأت|حدّثت|نُفِّذ)")
NEGATED = re.compile(r"(لم|لن|بدون|دون)\s*(يتم|أ?نشئ)?\s*")


def guard_input(text):
    """Screen text heading *into* the model — a question, or a document's body."""
    t = str(text or "")
    low = t.lower()
    hits = [name for pat, name in INJECTION_PATTERNS
            if re.search(pat, low, re.S)]
    return {
        "ok": not hits,
        "verdict": "block" if hits else "pass",
        "findings": hits,
        "chars": len(t),
    }


def guard_output(answer, session_id, has_pending_code=False):
    """Screen an answer heading *out*, against what the session was actually shown."""
    a = str(answer or "")
    findings = []

    for pat, name in SECRET_PATTERNS:
        if re.search(pat, a):
            findings.append({"kind": "secret_leak", "detail": name, "severity": "block"})

    ev = evidence_get(session_id) or {"codes": [], "ids": [], "sources": []}
    known = set(ev["codes"])
    cited = set(CODE_RE.findall(a))
    unsupported = sorted(cited - known)
    if unsupported and ev["sources"]:
        # Only meaningful once the session has been served something; before
        # that there is nothing to contradict.
        findings.append({"kind": "unsupported_record", "detail": unsupported,
                         "severity": "revise"})

    if re.search(UUID_ANY, a):
        findings.append({"kind": "raw_id_exposed", "detail": "uuid in answer",
                         "severity": "revise"})

    if CLAIMED_EXECUTION.search(NEGATED.sub(" NEG ", a)) and not has_pending_code:
        findings.append({"kind": "claimed_execution", "severity": "block",
                         "detail": "ادّعاء تنفيذ بلا رقم تأكيد"})

    block = any(f["severity"] == "block" for f in findings)
    return {
        "ok": not findings,
        "verdict": "block" if block else ("revise" if findings else "pass"),
        "findings": findings,
        "evidence": {"codes_known": len(known), "codes_cited": len(cited),
                     "sources": len(ev["sources"])},
    }


def acms_get_as(token, path):
    """Acms.get, but under the caller's token when one was forwarded, so a read
    through /q returns exactly what that caller would have seen directly."""
    req = urllib.request.Request(API + path)
    # No service-account fallback. With it, a call carrying no token at all was
    # served under the service account's permissions, which made every /q read
    # anonymous to anyone who could reach the port.
    if not token:
        raise RuntimeError("missing caller token")
    req.add_header("Authorization", "Bearer " + token)
    try:
        with urllib.request.urlopen(req, timeout=60) as r:
            raw = r.read().decode()
        return json.loads(raw) if raw else None
    except urllib.error.HTTPError as e:
        raise RuntimeError("%s -> %s %s" % (path, e.code, e.read().decode()[:200]))


def _project(rec, fields):
    out = {}
    for f in fields:
        if f in _FLAT:
            obj, key = _FLAT[f]
            src = rec.get(obj)
            v = src.get(key) if isinstance(src, dict) else None
        else:
            v = rec.get(f)
        # Dropping empties is a large part of the saving, and costs nothing:
        # a null serialises as wide as a value and tells the model less.
        # 0 and False are values, so they survive this test.
        if v is not None and v != "" and v != [] and v != {}:
            out[f] = v
    return out


def _iso(v):
    """Parse an ACMS timestamp. They arrive as ISO with a Z, which
    fromisoformat only learned to take in 3.11 — this box runs 3.10."""
    if not v:
        return None
    try:
        return datetime.fromisoformat(str(v).replace("Z", "+00:00"))
    except ValueError:
        return None


def _days_between(later, earlier):
    return int((later - earlier).total_seconds() // 86400)


def _enrich(resource, rows, token):
    """Attach computed fields and a facts summary. Returns (rows, facts).

    A failure to fetch the side data must not fail the read: the caller still
    gets its records, just without the derived numbers, and `facts` says so.
    """
    now = datetime.now(timezone.utc)
    facts = {"asOf": now.isoformat(timespec="seconds")}

    if resource == "accounts":
        try:
            raw = acms_get_as(token, "/activities")
            acts = (raw or {}).get("items") if isinstance(raw, dict) else (raw or [])
        except RuntimeError as e:
            facts["note"] = "تعذّر حساب زمن آخر تواصل: %s" % str(e)[:120]
            return rows, facts

        # Only a *completed* activity in the past counts as contact. A task that
        # is merely scheduled is not a contact, and an earlier version of this
        # calculation that took max(dueAt, completedAt) reported a last contact
        # dated in the future.
        last = {}
        for a in acts or []:
            acc = (a.get("account") or {}).get("legalName")
            done = _iso(a.get("completedAt"))
            if not acc or not done or done > now:
                continue
            if acc not in last or done > last[acc]:
                last[acc] = done

        never, gaps = [], []
        for r in rows:
            name = r.get("legalName")
            when = last.get(name)
            if when:
                r["lastContactAt"] = when.isoformat(timespec="seconds")
                r["daysSinceContact"] = _days_between(now, when)
                gaps.append((name, r["daysSinceContact"]))
            else:
                r["neverContacted"] = True
                never.append(name)
        facts["neverContacted"] = never
        facts["contactGapDays"] = dict(gaps)
        return rows, facts

    if resource == "activities":
        overdue = 0
        completed = 0
        for r in rows:
            done = _iso(r.get("completedAt"))
            due = _iso(r.get("dueAt"))
            if done:
                completed += 1
                continue
            if due and due < now:
                r["overdue"] = True
                r["daysOverdue"] = _days_between(now, due)
                overdue += 1
        facts["overdueCount"] = overdue
        facts["completedCount"] = completed
        facts["openCount"] = len(rows) - completed
        return rows, facts

    return rows, facts


def q_read(resource, params, token):
    spec = Q_SPECS.get(resource)
    if not spec:
        return 404, {"error": 'unknown resource "%s". available: %s'
                              % (resource, ", ".join(sorted(Q_SPECS)))}

    view = str(params.get("view") or "list").lower()
    if view not in spec["views"]:
        return 400, {"error": 'unknown view "%s". available: %s'
                              % (view, ", ".join(sorted(spec["views"])))}
    fields = list(spec["views"][view])

    # An explicit field list beats the view, for the rare question no preset covers.
    want = str(params.get("fields") or "").strip()
    if want:
        fields = [f.strip() for f in want.split(",") if f.strip()]

    try:
        limit = int(params.get("limit") or Q_LIMIT_DEFAULT)
    except (TypeError, ValueError):
        limit = Q_LIMIT_DEFAULT
    limit = max(1, min(limit, Q_LIMIT_MAX))

    qs = [(k, params[k]) for k in spec["passthrough"]
          if params.get(k) not in (None, "")]
    path = spec["path"] + ("?" + urllib.parse.urlencode(qs) if qs else "")

    if not token:
        # A missing caller token is a refusal, not an upstream fault.
        return 401, {"error": "missing caller token — /q requires the caller's ACMS bearer"}
    try:
        raw = acms_get_as(token, path)
    except RuntimeError as e:
        return 502, {"error": str(e)}

    items = raw.get("items") if isinstance(raw, dict) else raw
    if not isinstance(items, list):
        items = [raw] if raw else []

    for key in spec["local"]:
        v = params.get(key)
        if v in (None, ""):
            continue
        if key == "completed":
            done = str(v).lower() in ("1", "true", "yes")
            items = [r for r in items if bool(r.get("completedAt")) == done]
        else:
            items = [r for r in items
                     if str(r.get(key) or "").upper() == str(v).upper()]

    total = len(items)
    rows = [_project(r, fields) for r in items[:limit]]

    # The hard bound. `limit` caps the row count, but a row's width depends on
    # the view and on how much text each record happens to carry, so a row cap
    # alone does not cap the payload. Drop rows until the body fits, and report
    # what was dropped rather than shipping a silently short list.
    while len(rows) > 1 and len(json.dumps(rows, ensure_ascii=False)) > Q_MAX_CHARS:
        rows.pop()

    rows, facts = _enrich(resource, rows, token)

    out = {
        "items": rows,
        # Pre-computed so the answer quotes arithmetic instead of attempting it.
        "facts": facts,
        "returned": len(rows),
        "total": total,
        # Said outright, so an answer built on a slice can say it is a slice
        # instead of presenting it as the whole set.
        "truncated": total > len(rows),
        "view": view,
    }
    # Record what this session has now been shown, so an answer can be checked
    # against it rather than against the model's memory of it.
    evidence_record(params.get("session"), resource, out)

    if out["truncated"]:
        out["note"] = ("عُرض %d من %d سجلًا فقط. ضيّق البحث بفلتر "
                       "(status أو stage أو country أو search) للحصول على البقية."
                       % (len(rows), total))
    return 200, out


class Handler(BaseHTTPRequestHandler):
    server_version = "acms-reports/1.0"

    def log_message(self, fmt, *args):
        print("%s - %s" % (self.address_string(), fmt % args), flush=True)

    def _json(self, code, obj):
        body = json.dumps(obj, ensure_ascii=False).encode()
        self.send_response(code)
        self.send_header("Content-Type", "application/json; charset=utf-8")
        self.send_header("Content-Length", str(len(body)))
        self.end_headers()
        self.wfile.write(body)

    def _session(self):
        name = self.path.split("/pending/", 1)[1].split("?")[0]
        return name if SAFE_NAME.match(name) else None

    def do_GET(self):
        if not _svc_ok(self):
            return self._json(401, {"error": "unauthorized"})
        if self.path in ("/health", "/"):
            return self._json(200, {"status": "ok", "chrome": bool(CHROME),
                                    "reports": list(REPORTS), "formats": list(FORMATS),
                                    "pending": len(_pending),
                                    "cache": dict(CACHE_STATS)})
        if self.path.startswith("/q/"):
            parsed = urllib.parse.urlparse(self.path)
            resource = parsed.path[len("/q/"):].strip("/")
            params = {k: v[0] for k, v in
                      urllib.parse.parse_qs(parsed.query).items()}
            auth = self.headers.get("Authorization") or ""
            token = auth[7:].strip() if auth[:7].lower() == "bearer " else None
            code, body = q_read(resource, params, token)
            return self._json(code, body)
        if self.path.startswith("/telemetry/summary"):
            q = urllib.parse.parse_qs(urllib.parse.urlparse(self.path).query)
            try:
                days = int((q.get("days") or ["7"])[0])
            except ValueError:
                days = 7
            return self._json(200, telemetry_summary(max(1, min(days, 365))))
        if self.path.split("?")[0] in ("/login", "/login/"):
            body = LOGIN_PAGE.encode("utf-8")
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.send_header("Content-Length", str(len(body)))
            self.end_headers()
            self.wfile.write(body)
            return
        if self.path.startswith("/identity/"):
            name = self.path[len("/identity/"):].split("?")[0]
            if not SAFE_NAME.match(name):
                return self._json(400, {"error": "bad session id"})
            code, body = identity_resolve(name)
            return self._json(code, body)
        if self.path.startswith("/identity/"):
            name = self.path[len("/identity/"):].split("?")[0]
            return self._json(200, {"ok": True, "removed": identity_drop(name)})
        if self.path.startswith("/context/"):
            name = self.path[len("/context/"):].split("?")[0]
            if not SAFE_NAME.match(name):
                return self._json(400, {"error": "bad session id"})
            c = context_get(name)
            return self._json(200, {"found": bool(c), "context": c})
        if self.path.startswith("/evidence/"):
            name = self.path[len("/evidence/"):].split("?")[0]
            if not SAFE_NAME.match(name):
                return self._json(400, {"error": "bad session id"})
            e = evidence_get(name)
            return self._json(200, {"found": bool(e), "evidence": e})
        if self.path.startswith("/pending/"):
            sid = self._session()
            if not sid:
                return self._json(400, {"error": "bad session id"})
            p = pending_get(sid)
            # ?announce=1 hands the proposal over exactly once, so the chat can
            # append the confirmation block on the turn it was created and not
            # repeat it on every later message.
            if p and "announce=1" in self.path and not p.get("announced"):
                with _pending_lock:
                    p["announced"] = True
                return self._json(200, {"found": True, "fresh": True, "proposal": p})
            return self._json(200, {"found": bool(p), "fresh": False, "proposal": p})
        if self.path.startswith("/files/"):
            name = self.path[len("/files/"):].split("?")[0]
            # Only ever serve a plain generated filename from the output dir.
            if not SAFE_NAME.match(name):
                return self._json(400, {"error": "bad file name"})
            path = os.path.join(FILES, name)
            if not os.path.isfile(path):
                return self._json(404, {"error": "not found"})
            ctype = ("application/pdf" if name.endswith(".pdf") else
                     "application/vnd.openxmlformats-officedocument."
                     "presentationml.presentation")
            self.send_response(200)
            self.send_header("Content-Type", ctype)
            self.send_header("Content-Length", str(os.path.getsize(path)))
            self.send_header("Content-Disposition", 'inline; filename="%s"' % name)
            self.end_headers()
            with open(path, "rb") as fh:
                shutil.copyfileobj(fh, self.wfile)
            return
        return self._json(404, {"error": "not found"})

    def _body(self):
        n = int(self.headers.get("Content-Length") or 0)
        return json.loads(self.rfile.read(n).decode() or "{}")

    def do_PUT(self):
        if not _svc_ok(self):
            return self._json(401, {"error": "unauthorized"})
        if self.path.startswith("/context/"):
            name = self.path[len("/context/"):].split("?")[0]
            if not SAFE_NAME.match(name):
                return self._json(400, {"error": "bad session id"})
            return self._json(200, {"ok": True,
                                    "context": context_put(name, self._body())})
        if not self.path.startswith("/pending/"):
            return self._json(404, {"error": "not found"})
        sid = self._session()
        if not sid:
            return self._json(400, {"error": "bad session id"})
        try:
            pending_put(sid, self._body())
        except Exception as e:
            return self._json(400, {"error": "bad json: %s" % e})
        return self._json(200, {"ok": True})

    def do_DELETE(self):
        if not _svc_ok(self):
            return self._json(401, {"error": "unauthorized"})
        if self.path.startswith("/context/"):
            name = self.path[len("/context/"):].split("?")[0]
            with _context_lock:
                removed = _context.pop(name, None) is not None
            return self._json(200, {"ok": True, "removed": removed})
        if not self.path.startswith("/pending/"):
            return self._json(404, {"error": "not found"})
        sid = self._session()
        if not sid:
            return self._json(400, {"error": "bad session id"})
        return self._json(200, {"ok": True, "removed": bool(pending_pop(sid))})

    def do_POST(self):
        if not _svc_ok(self):
            return self._json(401, {"error": "unauthorized"})
        # Consume-on-confirm: returns the proposal and drops it in one step, so a
        # code can never be replayed even if two confirmations race.
        if self.path.startswith("/pending/") and self.path.rstrip("/").endswith("/claim"):
            name = self.path.split("/pending/", 1)[1].split("/claim")[0]
            if not SAFE_NAME.match(name):
                return self._json(400, {"error": "bad session id"})
            try:
                code = str(self._body().get("code", ""))
            except Exception as e:
                return self._json(400, {"error": "bad json: %s" % e})
            with _pending_lock:
                _prune_pending()
                p = _pending.get(name)
                if p and code and code == str(p.get("code")):
                    del _pending[name]
                    return self._json(200, {"claimed": True, "proposal": p})
            return self._json(200, {"claimed": False})

        if self.path.rstrip("/") == "/identity/start":
            b = self._body()
            code, body = identity_start(b.get("email"), b.get("password"))
            return self._json(code, body)
        if self.path.rstrip("/") == "/identity/bind":
            b = self._body()
            sid = str(b.get("sessionId") or "")
            if not SAFE_NAME.match(sid):
                return self._json(400, {"error": "bad session id"})
            code, body = identity_bind(b.get("code"), sid)
            return self._json(code, body)
        if self.path.rstrip("/") == "/guard/input":
            b = self._body()
            return self._json(200, guard_input(b.get("text")))
        if self.path.rstrip("/") == "/guard/output":
            b = self._body()
            return self._json(200, guard_output(
                b.get("answer"), str(b.get("sessionId") or ""),
                bool(b.get("hasPendingCode"))))
        if self.path.rstrip("/") in ("/cache/lookup", "/cache/store"):
            try:
                p = self._body()
            except Exception as e:
                return self._json(400, {"error": "bad json: %s" % e})
            user = str(p.get("user") or "").strip() or "svc"
            question = str(p.get("question") or "")
            if not question:
                return self._json(200, {"hit": False, "ok": False,
                                        "reason": "no-question"})
            if self.path.rstrip("/").endswith("lookup"):
                return self._json(200, cache_lookup(user, question))
            return self._json(200, cache_store(user, question,
                                               p.get("answer")))

        if self.path.rstrip("/") == "/propose":
            try:
                p = self._body()
            except Exception as e:
                return self._json(400, {"error": "bad json: %s" % e})
            sid = str(p.get("sessionId") or "").strip()
            if not sid or not SAFE_NAME.match(sid):
                return self._json(400, {"ok": False, "error": "bad sessionId"})
            params = p.get("params")
            if isinstance(params, str):
                try:
                    params = json.loads(params or "{}")
                except Exception as e:
                    return self._json(200, {"ok": False,
                                            "error": "params ليست JSON صالحًا: %s" % e})
            if not isinstance(params, dict):
                params = {}
            return self._json(200, propose(sid, str(p.get("action") or ""),
                                           params, str(p.get("summary") or "")))

        if self.path.rstrip("/") != "/generate":
            return self._json(404, {"error": "not found"})
        try:
            payload = self._body()
        except Exception as e:
            return self._json(400, {"error": "bad json: %s" % e})
        try:
            return self._json(200, generate(payload))
        except ValueError as e:
            return self._json(400, {"error": str(e)})
        except Exception as e:
            return self._json(500, {"error": "%s: %s" % (type(e).__name__, e)})


if __name__ == "__main__":
    print("acms-reports on :%d  chrome=%s  files=%s" % (PORT, bool(CHROME), FILES),
          flush=True)
    ThreadingHTTPServer((BIND, PORT), Handler).serve_forever()
